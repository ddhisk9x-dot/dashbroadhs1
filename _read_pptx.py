import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

ppt = Presentation(r"C:\Users\Duong Hieu\Desktop\deep dashbroad\bao cao\DeepDashboard_Final.pptx")
for i, slide in enumerate(ppt.slides):
    print(f"\n=== SLIDE {i+1} ===")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text.strip())
